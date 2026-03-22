"""PowerPoint report generation from stable analysis-run payloads."""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
TEAL = RGBColor(0x2E, 0xC4, 0xB6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF8)
MID_GRAY = RGBColor(0x8C, 0x9B, 0xB2)
ACCENT_ORANGE = RGBColor(0xFF, 0x6B, 0x35)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class PayloadReportBuilder:
    def __init__(self, payload: dict, charts: dict[str, bytes]):
        self.payload = payload
        self.charts = charts
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H

    def _blank_slide(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _add_rect(self, slide, left, top, width, height, fill_color: RGBColor):
        shape = slide.shapes.add_shape(1, left, top, width, height)
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.color.rgb = fill_color
        return shape

    def _add_textbox(
        self,
        slide,
        text: str,
        left,
        top,
        width,
        height,
        font_size: int = 12,
        bold: bool = False,
        color: RGBColor = DARK_BLUE,
        align: PP_ALIGN = PP_ALIGN.LEFT,
        word_wrap: bool = True,
        italic: bool = False,
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = "Calibri"
        return txBox

    def _add_image_bytes(self, slide, img_bytes: bytes, left, top, width, height=None):
        img_stream = io.BytesIO(img_bytes)
        if height:
            slide.shapes.add_picture(img_stream, left, top, width=width, height=height)
        else:
            slide.shapes.add_picture(img_stream, left, top, width=width)

    def _header_bar(self, slide, title: str, subtitle: str = ""):
        self._add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), DARK_BLUE)
        self._add_textbox(slide, title, Inches(0.4), Inches(0.12), Inches(10), Inches(0.55), font_size=22, bold=True, color=WHITE)
        if subtitle:
            self._add_textbox(slide, subtitle, Inches(0.4), Inches(0.65), Inches(10), Inches(0.35), font_size=11, color=TEAL)
        self._add_rect(slide, 0, Inches(1.1), SLIDE_W, Inches(0.04), TEAL)
        self._add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), DARK_BLUE)
        self._add_textbox(slide, "tundralis.com  ·  Confidential", Inches(0.3), Inches(7.22), Inches(5), Inches(0.25), font_size=7, color=MID_GRAY)
        self._add_textbox(slide, "Key Driver Analysis", Inches(9), Inches(7.22), Inches(4), Inches(0.25), font_size=7, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    def _slide_title(self):
        p = self.payload
        slide = self._blank_slide()
        self._add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, DARK_BLUE)
        self._add_rect(slide, 0, 0, Inches(0.12), SLIDE_H, TEAL)
        self._add_rect(slide, 0, Inches(6.8), SLIDE_W, Inches(0.7), TEAL)
        self._add_textbox(slide, "tundralis", Inches(0.4), Inches(0.3), Inches(4), Inches(0.5), font_size=14, bold=True, color=TEAL)
        self._add_textbox(slide, "tundralis.com", Inches(0.4), Inches(0.75), Inches(4), Inches(0.35), font_size=10, color=MID_GRAY)
        self._add_textbox(slide, "Key Driver Analysis", Inches(0.4), Inches(1.8), Inches(12), Inches(1.2), font_size=40, bold=True, color=WHITE)
        self._add_textbox(slide, f"Drivers of {p['outcome']['label']}", Inches(0.4), Inches(3.0), Inches(12), Inches(0.7), font_size=22, color=TEAL)
        self._add_textbox(
            slide,
            f"Sample size: {p['input_summary']['rows_modeled']:,} respondents  ·  Predictors analyzed: {len(p['input_summary']['predictor_columns'])}  ·  Model R² = {p['model_diagnostics']['r_squared']:.3f}",
            Inches(0.4), Inches(3.9), Inches(12), Inches(0.4), font_size=11, color=MID_GRAY,
        )
        self._add_textbox(slide, "CONFIDENTIAL  ·  FOR CLIENT USE ONLY", Inches(0.4), Inches(6.87), Inches(12), Inches(0.3), font_size=9, bold=True, color=DARK_BLUE)

    def _slide_exec_summary(self):
        p = self.payload
        slide = self._blank_slide()
        self._header_bar(slide, "Executive Summary")
        top = sorted(p["drivers"], key=lambda d: d["opportunity"]["rank"])[:3]
        bullets = []
        for d in top:
            bullets.append(f"• {d['driver_label']} — {d['classification']} with opportunity score {d['opportunity']['score']:.2f}")
        summary = "\n".join(bullets)
        summary += f"\n\nOutcome mean: {p['outcome']['summary']['mean']:.2f} / model fit R²={p['model_diagnostics']['r_squared']:.3f}."
        self._add_rect(slide, Inches(0.4), Inches(1.3), Inches(12.5), Inches(4.2), LIGHT_GRAY)
        self._add_textbox(slide, summary, Inches(0.7), Inches(1.5), Inches(12.0), Inches(3.8), font_size=14, color=DARK_BLUE)

    def _slide_chart(self, title: str, subtitle: str, chart_key: str, width=Inches(12.0), left=Inches(0.65), top=Inches(1.35)):
        slide = self._blank_slide()
        self._header_bar(slide, title, subtitle)
        if chart_key in self.charts:
            self._add_image_bytes(slide, self.charts[chart_key], left, top, width=width)
        return slide

    def _slide_opportunity_table(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Where to Focus First", "Opportunity ranking")
        drivers = sorted(self.payload["drivers"], key=lambda d: d["opportunity"]["rank"])[:8]
        headers = ["Rank", "Driver", "Opp.", "Imp.", "Impact", "Perf."]
        widths = [Inches(0.8), Inches(3.4), Inches(1.2), Inches(1.2), Inches(1.4), Inches(1.4)]
        starts = [Inches(0.45)]
        for w in widths[:-1]:
            starts.append(starts[-1] + w)
        y = Inches(1.35)
        self._add_rect(slide, Inches(0.4), y, Inches(10.0), Inches(0.38), DARK_BLUE)
        for h, x, w in zip(headers, starts, widths):
            self._add_textbox(slide, h, x, y + Inches(0.04), w, Inches(0.3), font_size=9, bold=True, color=WHITE)
        for i, d in enumerate(drivers):
            row_y = Inches(1.73) + i * Inches(0.45)
            self._add_rect(slide, Inches(0.4), row_y, Inches(10.0), Inches(0.42), LIGHT_GRAY if i % 2 == 0 else WHITE)
            vals = [
                str(d['opportunity']['rank']), d['driver_label'], f"{d['opportunity']['score']:.2f}",
                f"{d['importance']['share_of_explained_variance']:.1f}", f"{d['impact']['one_point_dv_change']:.2f}",
                f"{d['performance']['mean']:.2f}"
            ]
            for val, x, w in zip(vals, starts, widths):
                self._add_textbox(slide, val, x, row_y + Inches(0.04), w, Inches(0.3), font_size=9, color=DARK_BLUE)

    def _slide_confidence(self):
        p = self.payload
        slide = self._blank_slide()
        self._header_bar(slide, "Method Agreement / Confidence", "How much the methods agree and how sparse the data was")
        text = (
            f"Method agreement: {p['model_diagnostics']['method_agreement']}\n"
            f"Nonlinear signal: {p['model_diagnostics']['nonlinear_signal']}\n"
            f"Rows input: {p['input_summary']['rows_input']}\n"
            f"Rows with valid DV: {p['input_summary'].get('rows_with_valid_dv', 0)}\n"
            f"Rows with valid DV + any predictor: {p['input_summary'].get('rows_with_valid_dv_and_any_predictor', 0)}\n"
            f"Rows modeled: {p['input_summary']['rows_modeled']}"
        )
        self._add_rect(slide, Inches(0.4), Inches(1.4), Inches(5.8), Inches(2.4), LIGHT_GRAY)
        self._add_textbox(slide, text, Inches(0.7), Inches(1.7), Inches(5.2), Inches(1.9), font_size=14, color=DARK_BLUE)
        missing = self.payload['input_summary'].get('missingness', {}).get('by_variable', {})
        top_missing = sorted(missing.items(), key=lambda kv: kv[1].get('missing_rate', 0), reverse=True)[:8]
        self._add_textbox(slide, "Highest missingness variables", Inches(6.7), Inches(1.45), Inches(4), Inches(0.3), font_size=11, bold=True, color=DARK_BLUE)
        for i, (name, meta) in enumerate(top_missing):
            self._add_textbox(slide, f"• {name}: {meta['missing_rate']:.1%} missing", Inches(6.8), Inches(1.85 + i * 0.38), Inches(5.2), Inches(0.25), font_size=10, color=DARK_BLUE)

    def _slide_recommendations(self):
        slide = self._blank_slide()
        self._header_bar(slide, "Recommendations", "Prioritized actions based on KDA")
        recs = self.payload.get('recommendations', [])[:6]
        for i, rec in enumerate(recs):
            y = Inches(1.35) + i * Inches(0.92)
            self._add_rect(slide, Inches(0.35), y + Inches(0.08), Inches(0.4), Inches(0.4), TEAL)
            self._add_textbox(slide, str(i + 1), Inches(0.35), y + Inches(0.08), Inches(0.4), Inches(0.4), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            self._add_textbox(slide, rec['rationale'], Inches(0.9), y, Inches(12.0), Inches(0.8), font_size=11, color=DARK_BLUE)

    def build(self) -> Presentation:
        logger.info("Building PowerPoint report from analysis payload...")
        self._slide_title()
        self._slide_exec_summary()
        self._slide_chart("What Matters Most", "Importance ranking", "importance_bar", width=Inches(11.8), left=Inches(0.75), top=Inches(1.45))
        self._slide_chart("Priority Matrix", "Importance vs performance", "quadrant", width=Inches(8.8), left=Inches(0.5), top=Inches(1.25))
        self._slide_opportunity_table()
        self._slide_confidence()
        self._slide_recommendations()
        logger.info("Payload report complete: %d slides", len(self.prs.slides))
        return self.prs

    def save(self, path: str | Path) -> Path:
        path = Path(path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))
        logger.info("Saved payload report to %s", path)
        return path
